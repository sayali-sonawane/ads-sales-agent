"""
LangGraph agents for the LinkedIn Ads Sales Agent system.
This module defines the individual agent nodes that make up the sales pitch generation workflow.
Each agent is responsible for a specific part of the process: data analysis, visualization,
recommendation generation, and presentation creation.
"""

from typing import Dict, Any, List
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
from langchain_openai import ChatOpenAI
from langchain.prompts import ChatPromptTemplate
from langchain.schema import HumanMessage, SystemMessage

from sales_agent.graph_state import GraphState
from sales_agent.models import Recommendation, Industry, CustomerProfile
from sales_agent.data_generator import DataGenerator


class DataAnalysisAgent:
    """
    Agent responsible for analyzing customer performance data and comparing it with
    industry benchmarks and competitor performance. Generates insights for visualization
    and recommendation engines.
    """
    
    def __init__(self, llm: ChatOpenAI):
        """Initialize the data analysis agent with a language model."""
        self.llm = llm
        
    def analyze_performance(self, state: GraphState) -> GraphState:
        """
        Analyze customer historical performance and generate insights.
        
        Args:
            state: Current graph state containing customer and performance data
            
        Returns:
            Updated state with analysis results
        """
        try:
            historical_data = state["historical_data"]
            market_benchmarks = state["market_benchmarks"]
            competitor_data = state["competitor_data"]
            
            # Convert historical data to DataFrame for analysis
            df_data = []
            for hist in historical_data:
                df_data.append({
                    "period": hist.period,
                    "impressions": hist.metrics.impressions,
                    "clicks": hist.metrics.clicks,
                    "conversions": hist.metrics.conversions,
                    "ctr": hist.metrics.click_through_rate,
                    "conversion_rate": hist.metrics.conversion_rate,
                    "cpc": hist.metrics.cost_per_click,
                    "roas": hist.metrics.return_on_ad_spend,
                    "ad_spend": hist.ad_spend,
                    "revenue": hist.revenue_generated,
                    "engagement_rate": hist.metrics.engagement_rate
                })
            
            df = pd.DataFrame(df_data)
            
            # Calculate performance trends
            analysis_results = {}
            
            # Growth metrics
            if len(df) > 1:
                analysis_results["impressions_growth"] = (df["impressions"].iloc[-1] / df["impressions"].iloc[0] - 1) * 100
                analysis_results["revenue_growth"] = (df["revenue"].iloc[-1] / df["revenue"].iloc[0] - 1) * 100
                analysis_results["roas_trend"] = df["roas"].pct_change().mean() * 100
                analysis_results["ctr_trend"] = df["ctr"].pct_change().mean() * 100
            
            # Current performance vs benchmarks
            current_metrics = df.iloc[-1]
            if market_benchmarks:
                analysis_results["ctr_vs_benchmark"] = (current_metrics["ctr"] / market_benchmarks.avg_ctr - 1) * 100
                analysis_results["conversion_rate_vs_benchmark"] = (current_metrics["conversion_rate"] / market_benchmarks.avg_conversion_rate - 1) * 100
                analysis_results["roas_vs_benchmark"] = (current_metrics["roas"] / market_benchmarks.avg_roas - 1) * 100
                analysis_results["cpc_vs_benchmark"] = (current_metrics["cpc"] / market_benchmarks.avg_cost_per_click - 1) * 100
            
            # Competitor comparison
            if competitor_data:
                avg_competitor_engagement = sum(comp.avg_engagement_rate for comp in competitor_data) / len(competitor_data)
                analysis_results["engagement_vs_competitors"] = (current_metrics["engagement_rate"] / avg_competitor_engagement - 1) * 100
                
                total_competitor_spend = sum(comp.estimated_ad_spend for comp in competitor_data)
                analysis_results["spend_vs_competitors"] = (current_metrics["ad_spend"] / (total_competitor_spend / len(competitor_data)) - 1) * 100
            
            # Performance categories
            analysis_results["top_performing_metrics"] = []
            analysis_results["underperforming_metrics"] = []
            
            if market_benchmarks:
                if analysis_results.get("ctr_vs_benchmark", 0) > 10:
                    analysis_results["top_performing_metrics"].append("Click-through rate")
                elif analysis_results.get("ctr_vs_benchmark", 0) < -10:
                    analysis_results["underperforming_metrics"].append("Click-through rate")
                    
                if analysis_results.get("roas_vs_benchmark", 0) > 15:
                    analysis_results["top_performing_metrics"].append("Return on ad spend")
                elif analysis_results.get("roas_vs_benchmark", 0) < -15:
                    analysis_results["underperforming_metrics"].append("Return on ad spend")
            
            # Calculate seasonal patterns
            if len(df) >= 4:  # Need at least 4 quarters for seasonal analysis
                analysis_results["seasonal_patterns"] = {}
                for metric in ["impressions", "clicks", "conversions", "revenue"]:
                    quarterly_avg = df[metric].rolling(window=4).mean()
                    if not quarterly_avg.empty:
                        analysis_results["seasonal_patterns"][metric] = {
                            "peak_quarter": df.loc[quarterly_avg.idxmax(), "period"] if not quarterly_avg.isna().all() else "N/A",
                            "variance": quarterly_avg.std() / quarterly_avg.mean() if quarterly_avg.mean() > 0 else 0
                        }
            
            state["analysis_results"] = analysis_results
            state["current_step"] = "analysis_complete"
            
        except Exception as e:
            state["error_message"] = f"Data analysis failed: {str(e)}"
            state["current_step"] = "analysis_error"
        
        return state


class VisualizationAgent:
    """
    Agent responsible for creating charts and visualizations that will be included
    in the sales presentation. Generates bar charts, line graphs, and comparison charts
    based on the analysis results.
    """
    
    def __init__(self):
        """Initialize the visualization agent."""
        self.chart_style = {
            'figure.figsize': (12, 8),
            'axes.titlesize': 16,
            'axes.labelsize': 12,
            'xtick.labelsize': 10,
            'ytick.labelsize': 10,
            'legend.fontsize': 11
        }
        
    def create_visualizations(self, state: GraphState) -> GraphState:
        """
        Create visualizations based on analysis results and customer data.
        
        Args:
            state: Current graph state with analysis results
            
        Returns:
            Updated state with paths to generated visualization files
        """
        try:
            historical_data = state["historical_data"]
            analysis_results = state["analysis_results"]
            market_benchmarks = state["market_benchmarks"]
            competitor_data = state["competitor_data"]
            customer_profile = state["customer_profile"]
            
            visualization_paths = []
            
            # Set style for all plots
            plt.style.use('seaborn-v0_8')
            sns.set_palette("husl")
            
            # 1. Performance Trend Chart
            df_data = []
            for hist in historical_data:
                df_data.append({
                    "period": hist.period,
                    "impressions": hist.metrics.impressions,
                    "clicks": hist.metrics.clicks,
                    "conversions": hist.metrics.conversions,
                    "roas": hist.metrics.return_on_ad_spend,
                    "revenue": hist.revenue_generated,
                    "ad_spend": hist.ad_spend
                })
            
            df = pd.DataFrame(df_data)
            
            # Performance trends over time
            fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(15, 10))
            fig.suptitle(f'{customer_profile.company_name} - Performance Trends', fontsize=16, fontweight='bold')
            
            # ROAS trend
            ax1.plot(df["period"], df["roas"], marker='o', linewidth=2, markersize=6)
            ax1.set_title('Return on Ad Spend (ROAS)', fontweight='bold')
            ax1.set_ylabel('ROAS')
            ax1.tick_params(axis='x', rotation=45)
            ax1.grid(True, alpha=0.3)
            
            # Revenue vs Ad Spend
            ax2.bar(df["period"], df["revenue"], alpha=0.7, label='Revenue', color='green')
            ax2_twin = ax2.twinx()
            ax2_twin.plot(df["period"], df["ad_spend"], color='red', marker='s', label='Ad Spend', linewidth=2)
            ax2.set_title('Revenue vs Ad Spend', fontweight='bold')
            ax2.set_ylabel('Revenue ($)', color='green')
            ax2_twin.set_ylabel('Ad Spend ($)', color='red')
            ax2.tick_params(axis='x', rotation=45)
            ax2.legend(loc='upper left')
            ax2_twin.legend(loc='upper right')
            
            # Impressions and Clicks
            ax3.bar(df["period"], df["impressions"], alpha=0.6, label='Impressions', color='blue')
            ax3_twin = ax3.twinx()
            ax3_twin.bar(df["period"], df["clicks"], alpha=0.6, label='Clicks', color='orange', width=0.4)
            ax3.set_title('Impressions vs Clicks', fontweight='bold')
            ax3.set_ylabel('Impressions', color='blue')
            ax3_twin.set_ylabel('Clicks', color='orange')
            ax3.tick_params(axis='x', rotation=45)
            ax3.legend(loc='upper left')
            ax3_twin.legend(loc='upper right')
            
            # Conversions
            ax4.bar(df["period"], df["conversions"], alpha=0.7, color='purple')
            ax4.set_title('Conversions Over Time', fontweight='bold')
            ax4.set_ylabel('Conversions')
            ax4.tick_params(axis='x', rotation=45)
            ax4.grid(True, alpha=0.3)
            
            plt.tight_layout()
            trend_path = f"data/{customer_profile.company_name.replace(' ', '_')}_performance_trends.png"
            plt.savefig(trend_path, dpi=300, bbox_inches='tight')
            plt.close()
            visualization_paths.append(trend_path)
            
            # 2. Benchmark Comparison Chart
            if market_benchmarks:
                current_data = df.iloc[-1]
                
                fig, ax = plt.subplots(1, 1, figsize=(12, 8))
                
                metrics = ['CTR (%)', 'Conversion Rate (%)', 'ROAS', 'CPC ($)']
                
                # Calculate current performance
                customer_ctr = historical_data[-1].metrics.click_through_rate * 100
                customer_conv_rate = historical_data[-1].metrics.conversion_rate * 100
                customer_roas = historical_data[-1].metrics.return_on_ad_spend
                customer_cpc = historical_data[-1].metrics.cost_per_click
                
                customer_values = [customer_ctr, customer_conv_rate, customer_roas, customer_cpc]
                
                benchmark_values = [
                    market_benchmarks.avg_ctr * 100,
                    market_benchmarks.avg_conversion_rate * 100,
                    market_benchmarks.avg_roas,
                    market_benchmarks.avg_cost_per_click
                ]
                
                x = range(len(metrics))
                width = 0.35
                
                bars1 = ax.bar([i - width/2 for i in x], customer_values, width, 
                              label=f'{customer_profile.company_name}', alpha=0.8, color='steelblue')
                bars2 = ax.bar([i + width/2 for i in x], benchmark_values, width,
                              label=f'{customer_profile.industry.value} Benchmark', alpha=0.8, color='lightcoral')
                
                ax.set_title(f'Performance vs Industry Benchmarks - {customer_profile.industry.value}', 
                           fontsize=14, fontweight='bold')
                ax.set_ylabel('Value')
                ax.set_xticks(x)
                ax.set_xticklabels(metrics)
                ax.legend()
                ax.grid(True, alpha=0.3)
                
                # Add value labels on bars
                for bar in bars1:
                    height = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width()/2., height + height*0.01,
                           f'{height:.2f}', ha='center', va='bottom', fontweight='bold')
                
                for bar in bars2:
                    height = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width()/2., height + height*0.01,
                           f'{height:.2f}', ha='center', va='bottom', fontweight='bold')
                
                plt.tight_layout()
                benchmark_path = f"data/{customer_profile.company_name.replace(' ', '_')}_benchmark_comparison.png"
                plt.savefig(benchmark_path, dpi=300, bbox_inches='tight')
                plt.close()
                visualization_paths.append(benchmark_path)
            
            # 3. Competitor Analysis Chart
            if competitor_data:
                fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 6))
                
                # Market share pie chart
                competitor_names = [comp.competitor_name for comp in competitor_data]
                market_shares = [comp.market_share for comp in competitor_data]
                
                # Add customer's estimated market share
                total_competitor_share = sum(market_shares)
                customer_share = max(5, 100 - total_competitor_share - 20)  # Estimate customer share
                competitor_names.append(customer_profile.company_name)
                market_shares.append(customer_share)
                
                colors = sns.color_palette("husl", len(competitor_names))
                
                ax1.pie(market_shares, labels=competitor_names, autopct='%1.1f%%', 
                       colors=colors, startangle=90)
                ax1.set_title(f'Market Share - {customer_profile.industry.value}', fontweight='bold')
                
                # Ad spend comparison
                competitor_names_spend = [comp.competitor_name for comp in competitor_data]
                ad_spends = [comp.estimated_ad_spend for comp in competitor_data]
                
                # Add customer spend
                current_customer_spend = historical_data[-1].ad_spend
                competitor_names_spend.append(customer_profile.company_name + " (Current)")
                ad_spends.append(current_customer_spend)
                
                bars = ax2.bar(range(len(competitor_names_spend)), ad_spends, 
                              color=colors[:len(competitor_names_spend)])
                ax2.set_title('Estimated Monthly Ad Spend Comparison', fontweight='bold')
                ax2.set_ylabel('Ad Spend ($)')
                ax2.set_xticks(range(len(competitor_names_spend)))
                ax2.set_xticklabels(competitor_names_spend, rotation=45, ha='right')
                
                # Highlight customer bar
                bars[-1].set_color('darkred')
                bars[-1].set_alpha(0.9)
                
                # Add value labels
                for bar in bars:
                    height = bar.get_height()
                    ax2.text(bar.get_x() + bar.get_width()/2., height + height*0.01,
                           f'${height/1000:.0f}K', ha='center', va='bottom', fontweight='bold')
                
                plt.tight_layout()
                competitor_path = f"data/{customer_profile.company_name.replace(' ', '_')}_competitor_analysis.png"
                plt.savefig(competitor_path, dpi=300, bbox_inches='tight')
                plt.close()
                visualization_paths.append(competitor_path)
            
            state["visualizations"] = visualization_paths
            state["current_step"] = "visualizations_complete"
            
        except Exception as e:
            state["error_message"] = f"Visualization creation failed: {str(e)}"
            state["current_step"] = "visualization_error"
        
        return state


class RecommendationAgent:
    """
    Agent responsible for generating AI-powered recommendations based on the analysis
    results and market insights. Uses LLM to create strategic recommendations for
    improving LinkedIn advertising performance.
    """
    
    def __init__(self, llm: ChatOpenAI):
        """Initialize the recommendation agent with a language model."""
        self.llm = llm
        
        # Template for generating recommendations
        self.recommendation_prompt = ChatPromptTemplate.from_messages([
            ("system", """You are a LinkedIn advertising expert consultant. Based on the provided analysis data, 
            generate 3-5 strategic recommendations for improving the customer's LinkedIn advertising performance. 
            Each recommendation should include:
            1. A clear, actionable title
            2. Detailed description of the recommended action
            3. Expected impact and benefits
            4. Priority level (High/Medium/Low)
            5. Estimated ROI increase percentage
            6. Implementation timeline
            
            Focus on data-driven insights and LinkedIn-specific best practices. Consider industry benchmarks 
            and competitor analysis in your recommendations."""),
            ("human", """
            Customer Profile:
            - Company: {company_name}
            - Industry: {industry}
            - Company Size: {company_size}
            - Current Ad Spend: ${current_ad_spend:,.2f}
            - Primary Goals: {primary_goals}
            
            Performance Analysis:
            - Current ROAS: {current_roas:.2f}
            - CTR vs Benchmark: {ctr_vs_benchmark:+.1f}%
            - Conversion Rate vs Benchmark: {conversion_vs_benchmark:+.1f}%
            - ROAS vs Benchmark: {roas_vs_benchmark:+.1f}%
            - Revenue Growth: {revenue_growth:+.1f}%
            
            Top Performing Metrics: {top_performing}
            Underperforming Metrics: {underperforming}
            
            Competitor Insights:
            - Engagement vs Competitors: {engagement_vs_competitors:+.1f}%
            - Spend vs Competitors: {spend_vs_competitors:+.1f}%
            
            Please provide specific, actionable recommendations formatted as JSON with the following structure:
            {{
                "recommendations": [
                    {{
                        "title": "Recommendation title",
                        "description": "Detailed description",
                        "expected_impact": "Expected benefits and outcomes",
                        "priority": "High/Medium/Low",
                        "estimated_roi": 15.5,
                        "implementation_timeline": "2-4 weeks"
                    }}
                ]
            }}""")
        ])
    
    def generate_recommendations(self, state: GraphState) -> GraphState:
        """
        Generate strategic recommendations based on analysis results.
        
        Args:
            state: Current graph state with analysis results
            
        Returns:
            Updated state with generated recommendations
        """
        try:
            customer_profile = state["customer_profile"]
            analysis_results = state["analysis_results"]
            historical_data = state["historical_data"]
            
            # Get current performance metrics
            current_performance = historical_data[-1] if historical_data else None
            if not current_performance:
                raise ValueError("No historical performance data available")
            
            # If no LLM available, use fallback recommendations
            if not self.llm:
                recommendations = self._generate_fallback_recommendations(customer_profile, analysis_results)
                state["recommendations"] = recommendations
                state["current_step"] = "recommendations_complete"
                return state
            
            # Prepare data for the LLM prompt
            prompt_data = {
                "company_name": customer_profile.company_name,
                "industry": customer_profile.industry.value,
                "company_size": customer_profile.company_size,
                "current_ad_spend": customer_profile.current_ad_spend,
                "primary_goals": ", ".join(customer_profile.primary_goals),
                "current_roas": current_performance.metrics.return_on_ad_spend,
                "ctr_vs_benchmark": analysis_results.get("ctr_vs_benchmark", 0),
                "conversion_vs_benchmark": analysis_results.get("conversion_rate_vs_benchmark", 0),
                "roas_vs_benchmark": analysis_results.get("roas_vs_benchmark", 0),
                "revenue_growth": analysis_results.get("revenue_growth", 0),
                "top_performing": ", ".join(analysis_results.get("top_performing_metrics", ["None identified"])),
                "underperforming": ", ".join(analysis_results.get("underperforming_metrics", ["None identified"])),
                "engagement_vs_competitors": analysis_results.get("engagement_vs_competitors", 0),
                "spend_vs_competitors": analysis_results.get("spend_vs_competitors", 0)
            }
            
            # Generate recommendations using LLM
            formatted_prompt = self.recommendation_prompt.format(**prompt_data)
            response = self.llm.invoke(formatted_prompt)
            
            # Parse the response (in a real implementation, you'd want more robust JSON parsing)
            import json
            try:
                recommendations_data = json.loads(response.content)
                recommendations = []
                
                for rec_data in recommendations_data["recommendations"]:
                    recommendation = Recommendation(
                        title=rec_data["title"],
                        description=rec_data["description"],
                        expected_impact=rec_data["expected_impact"],
                        priority=rec_data["priority"],
                        estimated_roi=rec_data["estimated_roi"],
                        implementation_timeline=rec_data["implementation_timeline"]
                    )
                    recommendations.append(recommendation)
                
                state["recommendations"] = recommendations
                
            except json.JSONDecodeError:
                # Fallback to manual recommendations if LLM response parsing fails
                recommendations = self._generate_fallback_recommendations(customer_profile, analysis_results)
                state["recommendations"] = recommendations
            
            state["current_step"] = "recommendations_complete"
            
        except Exception as e:
            # Generate fallback recommendations in case of any errors
            recommendations = self._generate_fallback_recommendations(
                state["customer_profile"], 
                state["analysis_results"]
            )
            state["recommendations"] = recommendations
            state["current_step"] = "recommendations_complete"
        
        return state
    
    def _generate_fallback_recommendations(self, customer_profile: CustomerProfile, 
                                         analysis_results: Dict[str, Any]) -> List[Recommendation]:
        """
        Generate fallback recommendations when LLM-based generation fails.
        
        Args:
            customer_profile: Customer profile information
            analysis_results: Analysis results from data analysis
            
        Returns:
            List of fallback recommendations
        """
        recommendations = []
        
        # Recommendation 1: Audience Targeting Optimization
        recommendations.append(Recommendation(
            title="Optimize Audience Targeting for Higher CTR",
            description=f"Refine your LinkedIn audience targeting by leveraging job titles, skills, and company sizes that align with your {customer_profile.target_audience}. Implement lookalike audiences based on your highest-converting prospects.",
            expected_impact="Increase CTR by 25-40% and improve conversion quality by reaching more qualified prospects.",
            priority="High",
            estimated_roi=30.0,
            implementation_timeline="2-3 weeks"
        ))
        
        # Recommendation 2: Ad Creative Optimization
        recommendations.append(Recommendation(
            title="A/B Test Video and Carousel Ad Formats",
            description="LinkedIn video ads typically achieve 3x higher engagement. Test video testimonials, product demos, and carousel ads showcasing multiple value propositions to increase engagement rates.",
            expected_impact="Boost engagement rates by 45-60% and improve brand recall among target audience.",
            priority="High",
            estimated_roi=25.0,
            implementation_timeline="3-4 weeks"
        ))
        
        # Recommendation 3: Bid Strategy Optimization
        if analysis_results.get("cpc_vs_benchmark", 0) > 20:
            recommendations.append(Recommendation(
                title="Optimize Bid Strategy to Reduce Cost Per Click",
                description="Your CPC is above industry benchmark. Implement automated bidding with target CPA goals and consider dayparting to bid more aggressively during high-conversion hours.",
                expected_impact="Reduce CPC by 15-25% while maintaining conversion volume, improving overall ROAS.",
                priority="Medium",
                estimated_roi=20.0,
                implementation_timeline="1-2 weeks"
            ))
        
        # Recommendation 4: Landing Page Optimization
        if analysis_results.get("conversion_rate_vs_benchmark", 0) < -10:
            recommendations.append(Recommendation(
                title="Implement LinkedIn-Specific Landing Page Optimization",
                description="Create dedicated landing pages for LinkedIn traffic with B2B-focused messaging, social proof, and streamlined conversion forms. Implement LinkedIn Insight Tag for better attribution.",
                expected_impact="Increase conversion rates by 30-50% through better message matching and user experience.",
                priority="High",
                estimated_roi=35.0,
                implementation_timeline="4-6 weeks"
            ))
        
        # Recommendation 5: Budget Reallocation
        recommendations.append(Recommendation(
            title="Reallocate Budget to High-Performing Campaigns",
            description=f"Based on performance analysis, increase budget allocation to campaigns targeting {customer_profile.target_audience} during peak engagement periods. Reduce spend on underperforming segments.",
            expected_impact="Improve overall ROAS by 20-30% through more efficient budget distribution.",
            priority="Medium",
            estimated_roi=22.0,
            implementation_timeline="1 week"
        ))
        
        return recommendations[:4]  # Return top 4 recommendations


class PresentationAgent:
    """
    Agent responsible for creating the final PowerPoint presentation that combines
    all analysis results, visualizations, and recommendations into a cohesive
    sales pitch presentation.
    """
    
    def __init__(self):
        """Initialize the presentation agent."""
        pass
    
    def create_presentation(self, state: GraphState) -> GraphState:
        """
        Create a PowerPoint presentation with all the generated content.
        
        Args:
            state: Current graph state with all analysis results and visualizations
            
        Returns:
            Updated state with path to the generated presentation
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            import os
            
            customer_profile = state["customer_profile"]
            analysis_results = state["analysis_results"]
            recommendations = state["recommendations"]
            visualizations = state.get("visualizations", [])
            historical_data = state["historical_data"]
            
            # Create presentation
            prs = Presentation()
            
            # Slide 1: Title Slide
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = f"LinkedIn Advertising Strategy for {customer_profile.company_name}"
            subtitle.text = f"Performance Analysis & Strategic Recommendations\n{customer_profile.industry.value} Industry\nGenerated by LinkedIn Ads Sales Agent"
            
            # Slide 2: Executive Summary
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Executive Summary"
            
            # Get current performance for summary
            current_perf = historical_data[-1] if historical_data else None
            revenue_growth = analysis_results.get("revenue_growth", 0)
            roas_vs_benchmark = analysis_results.get("roas_vs_benchmark", 0)
            
            summary_text = f"""Key Findings:
            
            • Current monthly ad spend: ${current_perf.ad_spend:,.0f} if current_perf else 'N/A'
            • Revenue growth: {revenue_growth:+.1f}% over analysis period
            • ROAS vs industry benchmark: {roas_vs_benchmark:+.1f}%
            • {len(recommendations)} strategic recommendations identified
            
            Opportunity: Potential to increase ROI by 25-35% through targeted optimizations
            
            Primary Focus Areas:
            • Audience targeting refinement
            • Creative optimization for higher engagement
            • Budget reallocation to high-performing segments"""
            
            content.text = summary_text
            
            # Slide 3: Performance Trends
            if len(visualizations) > 0:
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = "Performance Trends Analysis"
                title_para.font.size = Pt(24)
                title_para.font.bold = True
                
                # Add chart image
                if os.path.exists(visualizations[0]):
                    slide.shapes.add_picture(visualizations[0], Inches(1), Inches(1.5), Inches(8), Inches(5))
            
            # Slide 4: Benchmark Comparison
            if len(visualizations) > 1:
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = f"Performance vs {customer_profile.industry.value} Benchmarks"
                title_para.font.size = Pt(24)
                title_para.font.bold = True
                
                # Add chart image
                if os.path.exists(visualizations[1]):
                    slide.shapes.add_picture(visualizations[1], Inches(1), Inches(1.5), Inches(8), Inches(5))
            
            # Slide 5: Competitor Analysis
            if len(visualizations) > 2:
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = "Competitive Landscape Analysis"
                title_para.font.size = Pt(24)
                title_para.font.bold = True
                
                # Add chart image
                if os.path.exists(visualizations[2]):
                    slide.shapes.add_picture(visualizations[2], Inches(1), Inches(1.5), Inches(8), Inches(5))
            
            # Slides 6+: Recommendations (one per slide)
            for i, rec in enumerate(recommendations):
                slide_layout = prs.slide_layouts[1]  # Title and content layout
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = f"Recommendation {i+1}: {rec.title}"
                
                rec_text = f"""Priority: {rec.priority}
                
                Description:
                {rec.description}
                
                Expected Impact:
                {rec.expected_impact}
                
                Estimated ROI Increase: +{rec.estimated_roi:.1f}%
                Implementation Timeline: {rec.implementation_timeline}"""
                
                content.text = rec_text
            
            # Final Slide: Next Steps
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Next Steps & Implementation"
            
            next_steps_text = """Immediate Actions (Next 30 Days):
            
            1. Schedule strategy review meeting to discuss recommendations
            2. Prioritize implementation based on business impact and resources
            3. Set up performance tracking for recommended optimizations
            4. Begin A/B testing for creative and audience optimizations
            
            LinkedIn Partnership Benefits:
            • Dedicated account management support
            • Access to beta features and advanced targeting options  
            • Regular performance reviews and optimization recommendations
            • Industry benchmark reports and competitive insights
            
            Contact your LinkedIn representative to begin implementation."""
            
            content.text = next_steps_text
            
            # Save presentation
            filename = f"{customer_profile.company_name.replace(' ', '_')}_LinkedIn_Strategy.pptx"
            presentation_path = f"data/{filename}"
            prs.save(presentation_path)
            
            state["presentation_path"] = presentation_path
            state["current_step"] = "presentation_complete"
            
        except Exception as e:
            state["error_message"] = f"Presentation creation failed: {str(e)}"
            state["current_step"] = "presentation_error"
        
        return state